import asyncio
import csv
import mimetypes
import os
from io import BufferedReader
from typing import Any, Dict, Optional

import docx2txt
import pptx
from fastapi import UploadFile
from loguru import logger
from pypdf import PdfReader
from pyzerox import zerox

from .schemas.document_schemas import DocumentMetadataSchema, DocumentSchema


async def get_document_from_file(
    file: UploadFile, metadata: DocumentMetadataSchema
) -> DocumentSchema:
    extracted_text = await extract_text_from_form_file(file)

    doc = DocumentSchema(text=extracted_text, metadata=metadata)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    try:
        with open(filepath, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(e)
        raise e

    return extracted_text


def extract_text_from_file(
    file: BufferedReader,
    mimetype: str,
    vision_config: Optional[Dict[str, Any]] = None,
) -> str:
    if vision_config and mimetype == "application/pdf":
        # Save to temporary file for vision model processing
        temp_file_path = "/tmp/temp_vision_file.pdf"
        with open(temp_file_path, "wb") as temp_file:
            temp_file.write(file.read())

        try:
            # Process with vision model
            extracted_text = asyncio.run(
                extract_text_with_vision_model(
                    file_path=temp_file_path,
                    model=vision_config.get("model", "gpt-4o-mini"),
                    api_key=vision_config.get("api_key"),
                    provider=vision_config.get("provider"),
                    system_prompt=vision_config.get("system_prompt"),
                )
            )
        finally:
            # Clean up temporary file
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)

        return extracted_text

    # Existing text extraction logic
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif mimetype == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    logger.info(f"mimetype: {mimetype}")
    logger.info(f"file.file: {file.file}")
    logger.info("file: ", file)

    file_stream = await file.read()

    temp_file_path = "/tmp/temp_file"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        logger.error(e)
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text


async def extract_text_with_vision_model(
    file_path: str,
    model: str = "gpt-4o-mini",
    api_key: Optional[str] = None,
    provider: Optional[str] = None,
    system_prompt: Optional[str] = None,
) -> str:
    """Extract text from a document using vision models via pyzerox."""
    kwargs: Dict[str, Any] = {}

    # Set up environment variables based on provider
    if provider == "openai" and api_key:
        os.environ["OPENAI_API_KEY"] = api_key
    elif provider == "azure" and api_key:
        os.environ["AZURE_API_KEY"] = api_key
    elif provider == "gemini" and api_key:
        os.environ["GEMINI_API_KEY"] = api_key
    elif provider == "anthropic" and api_key:
        os.environ["ANTHROPIC_API_KEY"] = api_key
    elif provider == "vertex_ai" and api_key:
        kwargs = {"vertex_credentials": api_key}

    try:
        # Process the document with zerox
        result = await zerox(
            file_path=file_path,
            model=model,
            output_dir="/tmp/zerox_output",  # Temporary output directory
            custom_system_prompt=system_prompt,
            cleanup=True,  # Clean up temporary files
            **kwargs,
        )
        return str(result)
    except Exception as e:
        logger.error(f"Error in vision model processing: {e}")
        raise e
