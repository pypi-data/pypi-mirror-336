"""
Reporting functionality for deduplication analysis and results.

This module provides capabilities to generate reports on deduplication results,
including metrics, visualizations, and detailed information about identified duplicates.
"""
from typing import Any, Dict, List, Optional, Union, Tuple
import os
import io
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import base64
from IPython.display import display, HTML, Markdown

from freamon.eda.report import (
    generate_header,
    generate_footer,
    fig_to_base64,
    add_table_styles
)


def generate_deduplication_report(
    results: Dict[str, Any],
    title: str = "Deduplication Analysis Report",
    format: str = "html",
    output_path: Optional[str] = None,
    include_pairs: bool = True,
    max_pairs: int = 100,
    theme: str = "cosmo"
) -> Union[str, None]:
    """
    Generate a report for deduplication analysis results.
    
    Parameters
    ----------
    results : Dict[str, Any]
        Dictionary containing deduplication results, should include:
        - 'duplicate_pairs': DataFrame or list of tuples with duplicate pairs
        - 'metrics': Dictionary with deduplication metrics
        - 'feature_importances': Optional dictionary of feature importances
        - 'model_info': Optional dictionary with model information
        - 'thresholds_data': Optional DataFrame with threshold evaluation data
    title : str, default="Deduplication Analysis Report"
        Title for the report
    format : str, default="html"
        Output format: "html", "markdown", "jupyter" (display in notebook)
    output_path : Optional[str], default=None
        Path to save the report (required for html and markdown formats)
    include_pairs : bool, default=True
        Whether to include sample duplicate pairs in the report
    max_pairs : int, default=100
        Maximum number of duplicate pairs to include
    theme : str, default="cosmo"
        Bootstrap theme for HTML report
        
    Returns
    -------
    Union[str, None]
        Report content as string (for "jupyter" format) or None if saved to file
    """
    if format not in ["html", "markdown", "jupyter"]:
        raise ValueError("Format must be one of: html, markdown, jupyter")
        
    if format != "jupyter" and output_path is None:
        raise ValueError("output_path is required for html and markdown formats")
        
    # Generate report content based on format
    if format == "html":
        report_content = _generate_html_report(
            results=results,
            title=title,
            include_pairs=include_pairs,
            max_pairs=max_pairs,
            theme=theme
        )
        
        # Save HTML report
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(report_content)
        return None
        
    elif format == "markdown":
        report_content = _generate_markdown_report(
            results=results,
            title=title,
            include_pairs=include_pairs,
            max_pairs=max_pairs
        )
        
        # Save Markdown report
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(report_content)
        return None
        
    else:  # jupyter
        report_content = _generate_markdown_report(
            results=results,
            title=title,
            include_pairs=include_pairs,
            max_pairs=max_pairs
        )
        
        # Display in Jupyter
        display(Markdown(report_content))
        return report_content


def _generate_html_report(
    results: Dict[str, Any],
    title: str,
    include_pairs: bool,
    max_pairs: int,
    theme: str
) -> str:
    """Generate an HTML report for deduplication results."""
    # Create header
    html = generate_header(title=title, theme=theme)
    
    # Add summary section
    html += f"<div class='container mt-5'>"
    html += f"<h1>{title}</h1>"
    html += f"<p class='text-muted'>Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>"
    
    # Add metrics section
    html += _create_metrics_section_html(results)
    
    # Add feature importances section if available
    if 'feature_importances' in results and results['feature_importances']:
        html += _create_feature_importances_section_html(results)
    
    # Add threshold evaluation section if available
    if 'thresholds_data' in results and results['thresholds_data'] is not None:
        html += _create_threshold_evaluation_section_html(results)
    
    # Add sample pairs section if requested
    if include_pairs and 'duplicate_pairs' in results and results['duplicate_pairs'] is not None:
        html += _create_sample_pairs_section_html(results, max_pairs)
    
    # Add model info section if available
    if 'model_info' in results and results['model_info']:
        html += _create_model_info_section_html(results)
    
    # Close container and add footer
    html += "</div>"
    html += generate_footer()
    
    return html


def _generate_markdown_report(
    results: Dict[str, Any],
    title: str,
    include_pairs: bool,
    max_pairs: int
) -> str:
    """Generate a Markdown report for deduplication results."""
    # Create header
    md = f"# {title}\n\n"
    md += f"*Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n"
    
    # Add metrics section
    md += _create_metrics_section_md(results)
    
    # Add feature importances section if available
    if 'feature_importances' in results and results['feature_importances']:
        md += _create_feature_importances_section_md(results)
    
    # Add threshold evaluation section if available
    if 'thresholds_data' in results and results['thresholds_data'] is not None:
        md += _create_threshold_evaluation_section_md(results)
    
    # Add sample pairs section if requested
    if include_pairs and 'duplicate_pairs' in results and results['duplicate_pairs'] is not None:
        md += _create_sample_pairs_section_md(results, max_pairs)
    
    # Add model info section if available
    if 'model_info' in results and results['model_info']:
        md += _create_model_info_section_md(results)
    
    return md


def _create_metrics_section_html(results: Dict[str, Any]) -> str:
    """Create HTML for metrics section."""
    if 'metrics' not in results or not results['metrics']:
        return ""
    
    metrics = results['metrics']
    
    html = "<div class='card mt-4 mb-4'>"
    html += "<div class='card-header bg-primary text-white'><h3>Deduplication Metrics</h3></div>"
    html += "<div class='card-body'>"
    html += "<div class='row'>"
    
    # Add general metrics
    html += "<div class='col-md-6'>"
    html += "<table class='table table-striped'>"
    html += "<thead><tr><th>Metric</th><th>Value</th></tr></thead>"
    html += "<tbody>"
    
    for metric, value in metrics.items():
        if isinstance(value, float):
            html += f"<tr><td>{metric.replace('_', ' ').title()}</td><td>{value:.4f}</td></tr>"
        else:
            html += f"<tr><td>{metric.replace('_', ' ').title()}</td><td>{value}</td></tr>"
    
    html += "</tbody></table>"
    html += "</div>"
    
    # Add visualization if available
    if 'metrics_chart' in results and results['metrics_chart'] is not None:
        fig = results['metrics_chart']
        img_str = fig_to_base64(fig)
        
        html += "<div class='col-md-6'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Metrics Chart'>"
        html += "</div>"
    
    html += "</div></div></div>"
    
    return html


def _create_metrics_section_md(results: Dict[str, Any]) -> str:
    """Create Markdown for metrics section."""
    if 'metrics' not in results or not results['metrics']:
        return ""
    
    metrics = results['metrics']
    
    md = "## Deduplication Metrics\n\n"
    md += "| Metric | Value |\n"
    md += "| ------ | ----- |\n"
    
    for metric, value in metrics.items():
        if isinstance(value, float):
            md += f"| {metric.replace('_', ' ').title()} | {value:.4f} |\n"
        else:
            md += f"| {metric.replace('_', ' ').title()} | {value} |\n"
    
    md += "\n"
    
    return md


def _create_feature_importances_section_html(results: Dict[str, Any]) -> str:
    """Create HTML for feature importances section."""
    feature_importances = results['feature_importances']
    
    # Convert to DataFrame for easier handling
    if isinstance(feature_importances, dict):
        fi_df = pd.DataFrame({
            'Feature': list(feature_importances.keys()),
            'Importance': list(feature_importances.values())
        }).sort_values(by='Importance', ascending=False)
    else:
        fi_df = feature_importances
    
    html = "<div class='card mt-4 mb-4'>"
    html += "<div class='card-header bg-success text-white'><h3>Feature Importances</h3></div>"
    html += "<div class='card-body'>"
    html += "<div class='row'>"
    
    # Add feature importance table
    html += "<div class='col-md-6'>"
    
    # Convert DataFrame to HTML table with styles
    table_html = fi_df.head(20).to_html(index=False, classes='table table-striped')
    table_html = add_table_styles(table_html)
    html += table_html
    
    html += "</div>"
    
    # Add visualization if available
    if 'feature_importance_chart' in results and results['feature_importance_chart'] is not None:
        fig = results['feature_importance_chart']
        img_str = fig_to_base64(fig)
        
        html += "<div class='col-md-6'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Feature Importances Chart'>"
        html += "</div>"
    # Otherwise create a simple bar chart
    elif not fi_df.empty:
        plt.figure(figsize=(10, 6))
        ax = sns.barplot(x='Importance', y='Feature', data=fi_df.head(15))
        ax.set_title('Top Features for Duplicate Detection')
        plt.tight_layout()
        
        img_str = fig_to_base64(plt.gcf())
        plt.close()
        
        html += "<div class='col-md-6'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Feature Importances Chart'>"
        html += "</div>"
    
    html += "</div></div></div>"
    
    return html


def _create_feature_importances_section_md(results: Dict[str, Any]) -> str:
    """Create Markdown for feature importances section."""
    feature_importances = results['feature_importances']
    
    # Convert to DataFrame for easier handling
    if isinstance(feature_importances, dict):
        fi_df = pd.DataFrame({
            'Feature': list(feature_importances.keys()),
            'Importance': list(feature_importances.values())
        }).sort_values(by='Importance', ascending=False)
    else:
        fi_df = feature_importances
    
    md = "## Feature Importances\n\n"
    
    # Format as markdown table
    md += "| Feature | Importance |\n"
    md += "| ------- | ---------- |\n"
    
    for _, row in fi_df.head(15).iterrows():
        md += f"| {row['Feature']} | {row['Importance']:.4f} |\n"
    
    md += "\n"
    
    return md


def _create_threshold_evaluation_section_html(results: Dict[str, Any]) -> str:
    """Create HTML for threshold evaluation section."""
    thresholds_data = results['thresholds_data']
    
    html = "<div class='card mt-4 mb-4'>"
    html += "<div class='card-header bg-info text-white'><h3>Threshold Evaluation</h3></div>"
    html += "<div class='card-body'>"
    
    # Add tabs for different visualizations
    html += "<ul class='nav nav-tabs' id='thresholdTabs' role='tablist'>"
    html += "<li class='nav-item'><a class='nav-link active' id='threshold-tab' data-toggle='tab' href='#threshold-chart' role='tab'>Threshold Chart</a></li>"
    html += "<li class='nav-item'><a class='nav-link' id='pr-tab' data-toggle='tab' href='#pr-curve' role='tab'>Precision-Recall Curve</a></li>"
    if 'cost_impact_chart' in results:
        html += "<li class='nav-item'><a class='nav-link' id='cost-tab' data-toggle='tab' href='#cost-impact' role='tab'>Cost Impact</a></li>"
    html += "</ul>"
    
    html += "<div class='tab-content p-3' id='thresholdTabContent'>"
    
    # Tab 1: Threshold Chart
    html += "<div class='tab-pane fade show active' id='threshold-chart' role='tabpanel'>"
    html += "<div class='row'>"
    
    # Add thresholds table
    html += "<div class='col-md-6'>"
    
    # Convert DataFrame to HTML table with styles
    table_html = thresholds_data.to_html(index=False, classes='table table-striped', float_format=lambda x: f"{x:.4f}")
    table_html = add_table_styles(table_html)
    html += table_html
    
    html += "</div>"
    
    # Add visualization if available
    if 'thresholds_chart' in results and results['thresholds_chart'] is not None:
        fig = results['thresholds_chart']
        img_str = fig_to_base64(fig)
        
        html += "<div class='col-md-6'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Threshold Evaluation Chart'>"
        html += "</div>"
    # Otherwise create a simple line chart
    elif not thresholds_data.empty:
        plt.figure(figsize=(10, 6))
        plt.plot(thresholds_data['threshold'], thresholds_data['precision'], 'b-', label='Precision')
        plt.plot(thresholds_data['threshold'], thresholds_data['recall'], 'r-', label='Recall')
        plt.plot(thresholds_data['threshold'], thresholds_data['f1'], 'g-', label='F1 Score')
        plt.xlabel('Probability Threshold')
        plt.ylabel('Score')
        plt.title('Precision, Recall, and F1 at Different Thresholds')
        plt.legend()
        plt.grid(True)
        plt.tight_layout()
        
        img_str = fig_to_base64(plt.gcf())
        plt.close()
        
        html += "<div class='col-md-6'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Threshold Evaluation Chart'>"
        html += "</div>"
    
    html += "</div>"
    html += "</div>"
    
    # Tab 2: Precision-Recall Curve
    html += "<div class='tab-pane fade' id='pr-curve' role='tabpanel'>"
    html += "<div class='row'>"
    
    # Add PR curve if available
    if 'pr_curve' in results and results['pr_curve'] is not None:
        fig = results['pr_curve']
        img_str = fig_to_base64(fig)
        
        html += "<div class='col-md-8'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Precision-Recall Curve'>"
        html += "</div>"
        
        html += "<div class='col-md-4'>"
        html += "<h5>Understanding the Precision-Recall Curve</h5>"
        html += "<p>This curve shows the trade-off between precision and recall at different probability thresholds:</p>"
        html += "<ul>"
        html += "<li><strong>Upper right corner</strong> is ideal (high precision and high recall)</li>"
        html += "<li><strong>Points to the right</strong> represent higher recall but may have lower precision</li>"
        html += "<li><strong>Points toward the top</strong> represent higher precision but may have lower recall</li>"
        html += "<li>The numbers on the curve represent the threshold values</li>"
        html += "</ul>"
        html += "<p>Choose a threshold that balances precision and recall based on your business needs.</p>"
        html += "</div>"
    else:
        # Create a simple PR curve
        html += "<div class='col-md-12'>"
        html += "<p class='alert alert-info'>Precision-Recall curve is not available for this dataset.</p>"
        html += "</div>"
    
    html += "</div>"
    html += "</div>"
    
    # Tab 3: Cost Impact (if available)
    if 'cost_impact_chart' in results:
        html += "<div class='tab-pane fade' id='cost-impact' role='tabpanel'>"
        html += "<div class='row'>"
        
        fig = results['cost_impact_chart']
        img_str = fig_to_base64(fig)
        
        html += "<div class='col-md-8'>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Cost Impact Chart'>"
        html += "</div>"
        
        html += "<div class='col-md-4'>"
        html += "<h5>Business Impact Analysis</h5>"
        html += "<p>This chart shows the estimated business cost at different threshold values:</p>"
        html += "<ul>"
        html += "<li>False positives cost: investigating non-duplicates</li>"
        html += "<li>False negatives cost: missing actual duplicates</li>"
        html += "<li>Total cost combines both with appropriate weights</li>"
        html += "</ul>"
        html += "<p>The optimal threshold minimizes the combined business impact.</p>"
        html += "</div>"
        
        html += "</div>"
        html += "</div>"
    
    html += "</div>"
    
    # Add precision-recall trade-off section
    html += "<div class='row mt-4'>"
    html += "<div class='col-md-12'>"
    html += "<h4>Precision-Recall Trade-off Analysis</h4>"
    html += "<p>The table below shows the trade-offs between precision and recall at different thresholds:</p>"
    
    # Create precision-recall trade-off table
    html += "<table class='table table-striped table-bordered'>"
    html += "<thead><tr><th>Threshold</th><th>Precision</th><th>Recall</th><th>True Positives</th><th>False Positives</th><th>False Negatives</th><th>Trade-off</th></tr></thead>"
    html += "<tbody>"
    
    # Add a row explaining the trade-off implications
    html += "<tr class='table-info'><td colspan='7'><strong>Trade-off Implications:</strong> "
    html += "Higher precision (higher threshold) means fewer false positives but more false negatives. "
    html += "Higher recall (lower threshold) means fewer false negatives but more false positives.</td></tr>"
    
    # Calculate TP, FP, FN estimates if available
    if 'metrics' in results and 'total_pairs_evaluated' in results['metrics']:
        total_pairs = results['metrics']['total_pairs_evaluated']
        true_duplicates = results['metrics'].get('true_duplicates', 0)
        
        for _, row in thresholds_data.iterrows():
            threshold = row['threshold']
            precision = row['precision']
            recall = row['recall']
            
            # Estimate TP, FP, FN
            tp = int(recall * true_duplicates)
            fp = int(tp * (1 - precision) / precision) if precision > 0 else "N/A"
            fn = int(true_duplicates - tp)
            
            # Determine trade-off description
            if threshold <= 0.3:
                trade_off = "High recall, lower precision: Catch most duplicates but with more false alarms"
            elif threshold <= 0.6:
                trade_off = "Balanced: Good compromise between precision and recall"
            else:
                trade_off = "High precision, lower recall: High confidence in duplicates found, but may miss some"
            
            html += f"<tr><td>{threshold:.2f}</td><td>{precision:.4f}</td><td>{recall:.4f}</td>"
            html += f"<td>{tp}</td><td>{fp}</td><td>{fn}</td><td>{trade_off}</td></tr>"
    else:
        # Without pair counts, just show qualitative trade-offs
        for _, row in thresholds_data.iterrows():
            threshold = row['threshold']
            precision = row['precision']
            recall = row['recall']
            
            # Determine trade-off description
            if threshold <= 0.3:
                trade_off = "High recall, lower precision: Catch most duplicates but with more false alarms"
            elif threshold <= 0.6:
                trade_off = "Balanced: Good compromise between precision and recall"
            else:
                trade_off = "High precision, lower recall: High confidence in duplicates found, but may miss some"
            
            html += f"<tr><td>{threshold:.2f}</td><td>{precision:.4f}</td><td>{recall:.4f}</td>"
            html += f"<td>-</td><td>-</td><td>-</td><td>{trade_off}</td></tr>"
    
    html += "</tbody></table>"
    
    # Add precision-recall curve if we have data
    if len(thresholds_data) >= 3:
        plt.figure(figsize=(10, 6))
        plt.plot(thresholds_data['recall'], thresholds_data['precision'], 'bo-')
        plt.xlabel('Recall')
        plt.ylabel('Precision')
        plt.title('Precision-Recall Curve')
        plt.grid(True)
        plt.xlim(0, 1.05)
        plt.ylim(0, 1.05)
        
        # Add threshold annotations
        for _, row in thresholds_data.iterrows():
            plt.annotate(f"{row['threshold']:.1f}", 
                        (row['recall'], row['precision']),
                        textcoords="offset points", 
                        xytext=(0,10), 
                        ha='center')
        
        plt.tight_layout()
        
        img_str = fig_to_base64(plt.gcf())
        plt.close()
        
        html += "<div class='mt-4'>"
        html += "<h4>Precision-Recall Curve</h4>"
        html += f"<img src='data:image/png;base64,{img_str}' class='img-fluid' alt='Precision-Recall Curve'>"
        html += "<p class='text-muted mt-2'>The numbers on the curve represent probability thresholds. "
        html += "The optimal point depends on your specific needs: prioritize precision for high-confidence matches, "
        html += "or recall to find more potential duplicates.</p>"
        html += "</div>"
    
    html += "</div>"
    html += "</div></div></div>"
    
    return html


def _create_threshold_evaluation_section_md(results: Dict[str, Any]) -> str:
    """Create Markdown for threshold evaluation section."""
    thresholds_data = results['thresholds_data']
    
    md = "## Threshold Evaluation\n\n"
    
    # Format as markdown table
    md += "| Threshold | Precision | Recall | F1 Score |\n"
    md += "| --------- | --------- | ------ | -------- |\n"
    
    for _, row in thresholds_data.iterrows():
        md += f"| {row['threshold']:.2f} | {row['precision']:.4f} | {row['recall']:.4f} | {row['f1']:.4f} |\n"
    
    md += "\n"
    
    # Add precision-recall trade-off section
    md += "### Precision-Recall Trade-offs\n\n"
    md += "The table below shows the trade-offs between precision and recall at different thresholds:\n\n"
    
    # Add trade-off explanation
    md += "> **Trade-off Implications:** Higher precision (higher threshold) means fewer false positives but more false negatives. "
    md += "Higher recall (lower threshold) means fewer false negatives but more false positives.\n\n"
    
    # Create trade-off table
    md += "| Threshold | Precision | Recall | True Positives | False Positives | False Negatives | Trade-off |\n"
    md += "| --------- | --------- | ------ | -------------- | --------------- | --------------- | --------- |\n"
    
    # Calculate TP, FP, FN estimates if available
    if 'metrics' in results and 'total_pairs_evaluated' in results['metrics']:
        total_pairs = results['metrics']['total_pairs_evaluated']
        true_duplicates = results['metrics'].get('true_duplicates', 0)
        
        for _, row in thresholds_data.iterrows():
            threshold = row['threshold']
            precision = row['precision']
            recall = row['recall']
            
            # Estimate TP, FP, FN
            tp = int(recall * true_duplicates)
            fp = int(tp * (1 - precision) / precision) if precision > 0 else "N/A"
            fn = int(true_duplicates - tp)
            
            # Determine trade-off description
            if threshold <= 0.3:
                trade_off = "High recall, lower precision: Find most duplicates with more false alarms"
            elif threshold <= 0.6:
                trade_off = "Balanced: Good compromise between precision and recall"
            else:
                trade_off = "High precision, lower recall: High confidence but may miss some"
            
            md += f"| {threshold:.2f} | {precision:.4f} | {recall:.4f} | {tp} | {fp} | {fn} | {trade_off} |\n"
    else:
        # Without pair counts, just show qualitative trade-offs
        for _, row in thresholds_data.iterrows():
            threshold = row['threshold']
            precision = row['precision']
            recall = row['recall']
            
            # Determine trade-off description
            if threshold <= 0.3:
                trade_off = "High recall, lower precision: Find most duplicates with more false alarms"
            elif threshold <= 0.6:
                trade_off = "Balanced: Good compromise between precision and recall"
            else:
                trade_off = "High precision, lower recall: High confidence but may miss some"
            
            md += f"| {threshold:.2f} | {precision:.4f} | {recall:.4f} | - | - | - | {trade_off} |\n"
    
    md += "\n"
    
    return md


def _create_sample_pairs_section_html(results: Dict[str, Any], max_pairs: int) -> str:
    """Create HTML for sample duplicate pairs section."""
    duplicate_pairs = results['duplicate_pairs']
    record_data = results.get('record_data', None)
    
    html = "<div class='card mt-4 mb-4'>"
    html += "<div class='card-header bg-warning'><h3>Sample Duplicate Pairs</h3></div>"
    html += "<div class='card-body'>"
    
    # Handle different formats of duplicate_pairs
    if isinstance(duplicate_pairs, pd.DataFrame):
        # Limit number of pairs
        pairs_df = duplicate_pairs.head(max_pairs)
        
        # Convert DataFrame to HTML table with styles
        table_html = pairs_df.to_html(classes='table table-striped', float_format=lambda x: f"{x:.4f}")
        table_html = add_table_styles(table_html)
        html += table_html
        
    elif isinstance(duplicate_pairs, list):
        # Create a simple list of pairs
        html += "<div class='table-responsive'>"
        html += "<table class='table table-striped'>"
        html += "<thead><tr><th>Index 1</th><th>Index 2</th></tr></thead>"
        html += "<tbody>"
        
        for i, pair in enumerate(duplicate_pairs[:max_pairs]):
            if isinstance(pair, tuple) and len(pair) == 2:
                html += f"<tr><td>{pair[0]}</td><td>{pair[1]}</td></tr>"
            
        html += "</tbody></table>"
        html += "</div>"
    
    # Add record data samples if available
    if record_data is not None and isinstance(record_data, list) and len(record_data) > 0:
        html += "<h4 class='mt-4'>Example Records</h4>"
        
        for i, record_pair in enumerate(record_data[:5]):  # Show up to 5 examples
            html += f"<div class='card mt-3'>"
            html += f"<div class='card-header'>Duplicate Pair {i+1}</div>"
            html += f"<div class='card-body'>"
            
            # Display record 1
            html += "<div class='row'>"
            html += "<div class='col-md-6'>"
            html += "<h5>Record 1</h5>"
            html += "<ul class='list-group'>"
            for field, value in record_pair[0].items():
                html += f"<li class='list-group-item'><strong>{field}:</strong> {value}</li>"
            html += "</ul>"
            html += "</div>"
            
            # Display record 2
            html += "<div class='col-md-6'>"
            html += "<h5>Record 2</h5>"
            html += "<ul class='list-group'>"
            for field, value in record_pair[1].items():
                html += f"<li class='list-group-item'><strong>{field}:</strong> {value}</li>"
            html += "</ul>"
            html += "</div>"
            html += "</div>"
            
            # Add probability if available
            if len(record_pair) > 2 and isinstance(record_pair[2], float):
                html += f"<div class='alert alert-info mt-3'>Duplicate probability: {record_pair[2]:.4f}</div>"
            
            html += "</div></div>"
    
    html += "</div></div>"
    
    return html


def _create_sample_pairs_section_md(results: Dict[str, Any], max_pairs: int) -> str:
    """Create Markdown for sample duplicate pairs section."""
    duplicate_pairs = results['duplicate_pairs']
    record_data = results.get('record_data', None)
    
    md = "## Sample Duplicate Pairs\n\n"
    
    # Handle different formats of duplicate_pairs
    if isinstance(duplicate_pairs, pd.DataFrame):
        # Convert DataFrame to markdown table
        pairs_df = duplicate_pairs.head(max_pairs)
        md += pairs_df.to_markdown(index=False)
        md += "\n\n"
        
    elif isinstance(duplicate_pairs, list):
        # Create a simple list of pairs
        md += "| Index 1 | Index 2 |\n"
        md += "| ------- | ------- |\n"
        
        for pair in duplicate_pairs[:max_pairs]:
            if isinstance(pair, tuple) and len(pair) == 2:
                md += f"| {pair[0]} | {pair[1]} |\n"
        
        md += "\n"
    
    # Add record data samples if available
    if record_data is not None and isinstance(record_data, list) and len(record_data) > 0:
        md += "### Example Records\n\n"
        
        for i, record_pair in enumerate(record_data[:5]):  # Show up to 5 examples
            md += f"#### Duplicate Pair {i+1}\n\n"
            
            # Display record 1
            md += "**Record 1:**\n\n"
            for field, value in record_pair[0].items():
                md += f"- **{field}:** {value}\n"
            md += "\n"
            
            # Display record 2
            md += "**Record 2:**\n\n"
            for field, value in record_pair[1].items():
                md += f"- **{field}:** {value}\n"
            md += "\n"
            
            # Add probability if available
            if len(record_pair) > 2 and isinstance(record_pair[2], float):
                md += f"**Duplicate probability:** {record_pair[2]:.4f}\n\n"
    
    return md


def _create_model_info_section_html(results: Dict[str, Any]) -> str:
    """Create HTML for model information section."""
    model_info = results['model_info']
    
    html = "<div class='card mt-4 mb-4'>"
    html += "<div class='card-header bg-secondary text-white'><h3>Model Information</h3></div>"
    html += "<div class='card-body'>"
    
    # Add model parameters and information
    html += "<table class='table table-striped'>"
    html += "<thead><tr><th>Parameter</th><th>Value</th></tr></thead>"
    html += "<tbody>"
    
    for param, value in model_info.items():
        if isinstance(value, dict):
            # For nested dictionaries, display nicely
            html += f"<tr><td>{param.replace('_', ' ').title()}</td><td><pre>{value}</pre></td></tr>"
        elif isinstance(value, list):
            # For lists, join with commas
            html += f"<tr><td>{param.replace('_', ' ').title()}</td><td>{', '.join(map(str, value))}</td></tr>"
        else:
            html += f"<tr><td>{param.replace('_', ' ').title()}</td><td>{value}</td></tr>"
    
    html += "</tbody></table>"
    html += "</div></div>"
    
    return html


def _create_model_info_section_md(results: Dict[str, Any]) -> str:
    """Create Markdown for model information section."""
    model_info = results['model_info']
    
    md = "## Model Information\n\n"
    
    md += "| Parameter | Value |\n"
    md += "| --------- | ----- |\n"
    
    for param, value in model_info.items():
        if isinstance(value, dict):
            # For nested dictionaries, use code formatting
            md += f"| {param.replace('_', ' ').title()} | ```{value}``` |\n"
        elif isinstance(value, list):
            # For lists, join with commas
            md += f"| {param.replace('_', ' ').title()} | {', '.join(map(str, value))} |\n"
        else:
            md += f"| {param.replace('_', ' ').title()} | {value} |\n"
    
    md += "\n"
    
    return md


def export_deduplication_report(
    results: Dict[str, Any],
    format: str = "excel",
    output_path: str = "deduplication_report",
    include_pairs: bool = True,
    max_pairs: int = 100
) -> None:
    """
    Export deduplication results to Excel or PowerPoint format.
    
    Parameters
    ----------
    results : Dict[str, Any]
        Dictionary containing deduplication results
    format : str, default="excel"
        Output format: "excel", "pptx", "html", "markdown"
    output_path : str, default="deduplication_report"
        Base path for the output file (without extension)
    include_pairs : bool, default=True
        Whether to include sample duplicate pairs
    max_pairs : int, default=100
        Maximum number of duplicate pairs to include
        
    Returns
    -------
    None
    """
    # Validate format
    if format not in ["excel", "pptx", "html", "markdown"]:
        raise ValueError("Format must be one of: excel, pptx, html, markdown")
    
    # Handle HTML and Markdown formats
    if format in ["html", "markdown"]:
        output_file = f"{output_path}.{format}"
        generate_deduplication_report(
            results=results,
            format=format,
            output_path=output_file,
            include_pairs=include_pairs,
            max_pairs=max_pairs
        )
        print(f"Report exported to {output_file}")
        return
    
    # Handle Excel format
    if format == "excel":
        _export_to_excel(
            results=results,
            output_path=f"{output_path}.xlsx",
            include_pairs=include_pairs,
            max_pairs=max_pairs
        )
        print(f"Report exported to {output_path}.xlsx")
        return
    
    # Handle PowerPoint format
    if format == "pptx":
        _export_to_powerpoint(
            results=results,
            output_path=f"{output_path}.pptx",
            include_pairs=include_pairs,
            max_pairs=max_pairs
        )
        print(f"Report exported to {output_path}.pptx")
        return


def _export_to_excel(
    results: Dict[str, Any],
    output_path: str,
    include_pairs: bool,
    max_pairs: int
) -> None:
    """Export deduplication results to Excel format."""
    from openpyxl import Workbook
    from openpyxl.utils.dataframe import dataframe_to_rows
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    
    wb = Workbook()
    
    # Create Summary sheet
    ws_summary = wb.active
    ws_summary.title = "Summary"
    
    # Add title
    ws_summary['A1'] = "Deduplication Analysis Report"
    ws_summary['A1'].font = Font(size=16, bold=True)
    ws_summary['A2'] = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    ws_summary['A2'].font = Font(italic=True)
    
    # Add metrics
    row = 4
    if 'metrics' in results and results['metrics']:
        ws_summary[f'A{row}'] = "Deduplication Metrics"
        ws_summary[f'A{row}'].font = Font(size=14, bold=True)
        row += 1
        
        metrics = results['metrics']
        ws_summary[f'A{row}'] = "Metric"
        ws_summary[f'B{row}'] = "Value"
        ws_summary[f'A{row}'].font = Font(bold=True)
        ws_summary[f'B{row}'].font = Font(bold=True)
        row += 1
        
        for metric, value in metrics.items():
            ws_summary[f'A{row}'] = metric.replace('_', ' ').title()
            if isinstance(value, float):
                ws_summary[f'B{row}'] = f"{value:.4f}"
            else:
                ws_summary[f'B{row}'] = str(value)
            row += 2
    
    # Create Feature Importances sheet
    if 'feature_importances' in results and results['feature_importances']:
        ws_fi = wb.create_sheet("Feature Importances")
        
        feature_importances = results['feature_importances']
        
        # Convert to DataFrame for easier handling
        if isinstance(feature_importances, dict):
            fi_df = pd.DataFrame({
                'Feature': list(feature_importances.keys()),
                'Importance': list(feature_importances.values())
            }).sort_values(by='Importance', ascending=False)
        else:
            fi_df = feature_importances
        
        # Add title
        ws_fi['A1'] = "Feature Importances for Duplicate Detection"
        ws_fi['A1'].font = Font(size=14, bold=True)
        
        # Add data
        for r_idx, row in enumerate(dataframe_to_rows(fi_df, index=False, header=True), 3):
            for c_idx, value in enumerate(row, 1):
                cell = ws_fi.cell(row=r_idx, column=c_idx, value=value)
                if r_idx == 3:  # Header row
                    cell.font = Font(bold=True)
        
        # Adjust column widths
        for column in ws_fi.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = (max_length + 2)
            ws_fi.column_dimensions[column_letter].width = adjusted_width
    
    # Create Threshold Evaluation sheet
    if 'thresholds_data' in results and results['thresholds_data'] is not None:
        ws_th = wb.create_sheet("Threshold Evaluation")
        
        thresholds_data = results['thresholds_data']
        
        # Add title
        ws_th['A1'] = "Threshold Evaluation Results"
        ws_th['A1'].font = Font(size=14, bold=True)
        
        # Add data
        for r_idx, row in enumerate(dataframe_to_rows(thresholds_data, index=False, header=True), 3):
            for c_idx, value in enumerate(row, 1):
                cell = ws_th.cell(row=r_idx, column=c_idx, value=value)
                if r_idx == 3:  # Header row
                    cell.font = Font(bold=True)
        
        # Adjust column widths
        for column in ws_th.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = (max_length + 2)
            ws_th.column_dimensions[column_letter].width = adjusted_width
    
    # Create Duplicate Pairs sheet
    if include_pairs and 'duplicate_pairs' in results and results['duplicate_pairs'] is not None:
        ws_pairs = wb.create_sheet("Duplicate Pairs")
        
        duplicate_pairs = results['duplicate_pairs']
        
        # Add title
        ws_pairs['A1'] = "Sample Duplicate Pairs"
        ws_pairs['A1'].font = Font(size=14, bold=True)
        
        # Handle different formats of duplicate_pairs
        if isinstance(duplicate_pairs, pd.DataFrame):
            # Add data
            pairs_df = duplicate_pairs.head(max_pairs)
            for r_idx, row in enumerate(dataframe_to_rows(pairs_df, index=False, header=True), 3):
                for c_idx, value in enumerate(row, 1):
                    cell = ws_pairs.cell(row=r_idx, column=c_idx, value=value)
                    if r_idx == 3:  # Header row
                        cell.font = Font(bold=True)
        
        elif isinstance(duplicate_pairs, list):
            # Create headers
            ws_pairs['A3'] = "Index 1"
            ws_pairs['B3'] = "Index 2"
            ws_pairs['A3'].font = Font(bold=True)
            ws_pairs['B3'].font = Font(bold=True)
            
            # Add data
            for i, pair in enumerate(duplicate_pairs[:max_pairs], 4):
                if isinstance(pair, tuple) and len(pair) == 2:
                    ws_pairs[f'A{i}'] = pair[0]
                    ws_pairs[f'B{i}'] = pair[1]
        
        # Adjust column widths
        for column in ws_pairs.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = (max_length + 2)
            ws_pairs.column_dimensions[column_letter].width = adjusted_width
    
    # Create Model Info sheet
    if 'model_info' in results and results['model_info']:
        ws_model = wb.create_sheet("Model Information")
        
        model_info = results['model_info']
        
        # Add title
        ws_model['A1'] = "Model Information"
        ws_model['A1'].font = Font(size=14, bold=True)
        
        # Add headers
        ws_model['A3'] = "Parameter"
        ws_model['B3'] = "Value"
        ws_model['A3'].font = Font(bold=True)
        ws_model['B3'].font = Font(bold=True)
        
        # Add data
        row = 4
        for param, value in model_info.items():
            ws_model[f'A{row}'] = param.replace('_', ' ').title()
            if isinstance(value, dict) or isinstance(value, list):
                ws_model[f'B{row}'] = str(value)
            else:
                ws_model[f'B{row}'] = value
            row += 1
        
        # Adjust column widths
        for column in ws_model.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 100)  # Cap at 100
            ws_model.column_dimensions[column_letter].width = adjusted_width
    
    # Save workbook
    wb.save(output_path)


def _export_to_powerpoint(
    results: Dict[str, Any],
    output_path: str,
    include_pairs: bool,
    max_pairs: int
) -> None:
    """Export deduplication results to PowerPoint format."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("python-pptx package is required for PowerPoint export. Install with: pip install python-pptx")
    
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Deduplication Analysis Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # Create metrics slide
    if 'metrics' in results and results['metrics']:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Deduplication Metrics"
        
        metrics = results['metrics']
        
        # Add content placeholder
        content = slide.placeholders[1]
        tf = content.text_frame
        
        for metric, value in metrics.items():
            p = tf.add_paragraph()
            if isinstance(value, float):
                p.text = f"{metric.replace('_', ' ').title()}: {value:.4f}"
            else:
                p.text = f"{metric.replace('_', ' ').title()}: {value}"
            p.font.size = Pt(18)
    
    # Create Feature Importances slide
    if 'feature_importances' in results and results['feature_importances']:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Feature Importances"
        
        content = slide.placeholders[1]
        left = content.left
        top = content.top
        width = content.width
        height = content.height
        
        # Add chart if available
        if 'feature_importance_chart' in results and results['feature_importance_chart'] is not None:
            fig = results['feature_importance_chart']
            # Save temporary image file
            tmp_img = "temp_fi_chart.png"
            fig.savefig(tmp_img, dpi=300, bbox_inches='tight')
            
            # Add image to slide
            slide.shapes.add_picture(tmp_img, left, top, width, height)
            
            # Remove temporary file
            os.remove(tmp_img)
        else:
            # Add text content instead
            feature_importances = results['feature_importances']
            
            # Convert to DataFrame for easier handling
            if isinstance(feature_importances, dict):
                fi_df = pd.DataFrame({
                    'Feature': list(feature_importances.keys()),
                    'Importance': list(feature_importances.values())
                }).sort_values(by='Importance', ascending=False)
            else:
                fi_df = feature_importances
            
            # Add table with top features
            top_features = fi_df.head(10)
            rows, cols = len(top_features) + 1, 2  # +1 for header
            
            table = slide.shapes.add_table(rows, cols, left, top, width * 0.8, height * 0.8).table
            
            # Set headers
            table.cell(0, 0).text = "Feature"
            table.cell(0, 1).text = "Importance"
            
            # Make headers bold
            for cell in table.rows[0].cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(14)
            
            # Fill data
            for i, (_, row) in enumerate(top_features.iterrows(), 1):
                table.cell(i, 0).text = str(row['Feature'])
                table.cell(i, 1).text = f"{row['Importance']:.4f}"
    
    # Create Threshold Evaluation slide
    if 'thresholds_data' in results and results['thresholds_data'] is not None:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Threshold Evaluation"
        
        content = slide.placeholders[1]
        left = content.left
        top = content.top
        width = content.width
        height = content.height
        
        # Add chart if available
        if 'thresholds_chart' in results and results['thresholds_chart'] is not None:
            fig = results['thresholds_chart']
            # Save temporary image file
            tmp_img = "temp_th_chart.png"
            fig.savefig(tmp_img, dpi=300, bbox_inches='tight')
            
            # Add image to slide
            slide.shapes.add_picture(tmp_img, left, top, width, height)
            
            # Remove temporary file
            os.remove(tmp_img)
        else:
            # Add text content instead
            thresholds_data = results['thresholds_data']
            
            # Create thresholds chart
            plt.figure(figsize=(10, 6))
            plt.plot(thresholds_data['threshold'], thresholds_data['precision'], 'b-', label='Precision')
            plt.plot(thresholds_data['threshold'], thresholds_data['recall'], 'r-', label='Recall')
            plt.plot(thresholds_data['threshold'], thresholds_data['f1'], 'g-', label='F1 Score')
            plt.xlabel('Probability Threshold')
            plt.ylabel('Score')
            plt.title('Precision, Recall, and F1 at Different Thresholds')
            plt.legend()
            plt.grid(True)
            plt.tight_layout()
            
            # Save temporary image file
            tmp_img = "temp_th_chart.png"
            plt.savefig(tmp_img, dpi=300, bbox_inches='tight')
            plt.close()
            
            # Add image to slide
            slide.shapes.add_picture(tmp_img, left, top, width, height)
            
            # Remove temporary file
            os.remove(tmp_img)
    
    # Create Sample Pairs slide
    if include_pairs and 'duplicate_pairs' in results and results['duplicate_pairs'] is not None and 'record_data' in results:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Sample Duplicate Pairs"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        record_data = results['record_data']
        
        if record_data is not None and isinstance(record_data, list) and len(record_data) > 0:
            # Show example record pairs
            for i, record_pair in enumerate(record_data[:3]):  # Show top 3 examples
                p = tf.add_paragraph()
                p.text = f"Pair {i+1}:"
                p.font.bold = True
                p.font.size = Pt(16)
                
                # Record 1
                p = tf.add_paragraph()
                p.text = "Record 1:"
                p.font.bold = True
                p.font.size = Pt(14)
                
                for field, value in list(record_pair[0].items())[:5]:  # Show top 5 fields
                    p = tf.add_paragraph()
                    p.text = f"  {field}: {value}"
                    p.font.size = Pt(12)
                
                # Record 2
                p = tf.add_paragraph()
                p.text = "Record 2:"
                p.font.bold = True
                p.font.size = Pt(14)
                
                for field, value in list(record_pair[1].items())[:5]:  # Show top 5 fields
                    p = tf.add_paragraph()
                    p.text = f"  {field}: {value}"
                    p.font.size = Pt(12)
                
                # Add probability if available
                if len(record_pair) > 2 and isinstance(record_pair[2], float):
                    p = tf.add_paragraph()
                    p.text = f"Duplicate probability: {record_pair[2]:.4f}"
                    p.font.size = Pt(14)
                    p.font.italic = True
                
                # Add space between pairs
                if i < min(3, len(record_data) - 1):
                    p = tf.add_paragraph()
                    p.text = "—" * 30
    
    # Create Model Info slide
    if 'model_info' in results and results['model_info']:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Model Information"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        model_info = results['model_info']
        
        for param, value in model_info.items():
            p = tf.add_paragraph()
            p.font.size = Pt(14)
            
            param_name = param.replace('_', ' ').title()
            if isinstance(value, dict):
                p.text = f"{param_name}: {str(value)[:100]}..."
            elif isinstance(value, list):
                p.text = f"{param_name}: {', '.join(map(str, value[:5]))}"
                if len(value) > 5:
                    p.text += "..."
            else:
                p.text = f"{param_name}: {value}"
    
    # Save presentation
    prs.save(output_path)


def prepare_duplicate_report_data(
    model,
    test_df: pd.DataFrame,
    true_duplicate_pairs: List[Tuple[int, int]],
    threshold_range: Optional[List[float]] = None,
    include_record_data: bool = True,
    max_record_samples: int = 5,
    sample_key_fields: Optional[List[str]] = None
) -> Dict[str, Any]:
    """
    Prepare data for deduplication report generation.
    
    Parameters
    ----------
    model : SupervisedDeduplicationModel
        The trained deduplication model
    test_df : pd.DataFrame
        Test dataframe for evaluation
    true_duplicate_pairs : List[Tuple[int, int]]
        List of known duplicate pairs in the test data
    threshold_range : Optional[List[float]], default=None
        List of thresholds to evaluate. If None, defaults to [0.1, 0.3, 0.5, 0.7, 0.9]
    include_record_data : bool, default=True
        Whether to include sample record data
    max_record_samples : int, default=5
        Maximum number of record samples to include
    sample_key_fields : Optional[List[str]], default=None
        List of key fields to include in record samples. If None, includes all fields.
        
    Returns
    -------
    Dict[str, Any]
        Dictionary containing organized data for report generation
    """
    # Set default threshold range if not provided
    if threshold_range is None:
        threshold_range = [0.1, 0.3, 0.5, 0.7, 0.9]
    
    results = {}
    
    # Evaluate model at different thresholds
    threshold_results = []
    for threshold in threshold_range:
        metrics = model.evaluate(
            df=test_df,
            true_duplicate_pairs=true_duplicate_pairs,
            threshold=threshold
        )
        threshold_results.append({
            'threshold': threshold,
            **metrics
        })
    
    # Create thresholds dataframe
    results['thresholds_data'] = pd.DataFrame(threshold_results)
    
    # Find "best" threshold based on F1 score
    best_result = max(threshold_results, key=lambda x: x['f1'])
    best_threshold = best_result['threshold']
    
    # Add metrics using best threshold
    results['metrics'] = {
        'precision': best_result['precision'],
        'recall': best_result['recall'],
        'f1_score': best_result['f1'],
        'auc': best_result['auc'],
        'best_threshold': best_threshold,
        'total_pairs_evaluated': len(true_duplicate_pairs),
        'true_duplicates': len(true_duplicate_pairs)
    }
    
    # Get feature importances
    if hasattr(model, 'get_feature_importances'):
        results['feature_importances'] = model.get_feature_importances()
    
    # Get duplicate predictions using best threshold
    predictions = model.find_duplicates(
        df=test_df,
        threshold=best_threshold,
        return_probabilities=True
    )
    results['duplicate_pairs'] = predictions
    
    # Add model information
    model_info = {
        'model_type': model.model_type,
        'best_threshold': best_threshold,
    }
    
    if hasattr(model, 'key_features'):
        model_info['key_features'] = model.key_features
    
    if hasattr(model, 'date_features'):
        model_info['date_features'] = model.date_features
    
    if hasattr(model, 'model_params'):
        model_info['model_params'] = model.model_params
    
    results['model_info'] = model_info
    
    # Include sample record data if requested
    if include_record_data:
        record_data = []
        
        # Get sample pairs
        sample_pairs = predictions.head(max_record_samples).values
        
        for pair in sample_pairs:
            idx1, idx2 = int(pair[0]), int(pair[1])
            prob = float(pair[2])
            
            # Get record data
            record1 = test_df.iloc[idx1].to_dict()
            record2 = test_df.iloc[idx2].to_dict()
            
            # Filter fields if sample_key_fields provided
            if sample_key_fields:
                record1 = {k: v for k, v in record1.items() if k in sample_key_fields}
                record2 = {k: v for k, v in record2.items() if k in sample_key_fields}
            
            record_data.append((record1, record2, prob))
        
        results['record_data'] = record_data
    
    # Create threshold evaluation chart
    plt.figure(figsize=(10, 6))
    plt.plot(results['thresholds_data']['threshold'], results['thresholds_data']['precision'], 'b-', label='Precision')
    plt.plot(results['thresholds_data']['threshold'], results['thresholds_data']['recall'], 'r-', label='Recall')
    plt.plot(results['thresholds_data']['threshold'], results['thresholds_data']['f1'], 'g-', label='F1 Score')
    plt.xlabel('Probability Threshold')
    plt.ylabel('Score')
    plt.title('Precision, Recall, and F1 at Different Thresholds')
    plt.legend()
    plt.grid(True)
    plt.tight_layout()
    results['thresholds_chart'] = plt.gcf()
    
    # Create precision-recall curve
    plt.figure(figsize=(10, 6))
    plt.plot(results['thresholds_data']['recall'], results['thresholds_data']['precision'], 'bo-')
    plt.xlabel('Recall')
    plt.ylabel('Precision')
    plt.title('Precision-Recall Curve')
    plt.grid(True)
    plt.xlim(0, 1.05)
    plt.ylim(0, 1.05)
    
    # Add threshold annotations
    for _, row in results['thresholds_data'].iterrows():
        plt.annotate(f"{row['threshold']:.1f}", 
                    (row['recall'], row['precision']),
                    textcoords="offset points", 
                    xytext=(0,10), 
                    ha='center')
    
    plt.tight_layout()
    results['pr_curve'] = plt.gcf()
    
    # Create feature importance chart if available
    if 'feature_importances' in results and results['feature_importances']:
        feature_importances = results['feature_importances']
        
        # Convert to DataFrame for easier plotting
        if isinstance(feature_importances, dict):
            fi_df = pd.DataFrame({
                'Feature': list(feature_importances.keys()),
                'Importance': list(feature_importances.values())
            }).sort_values(by='Importance', ascending=False)
        else:
            fi_df = feature_importances
        
        plt.figure(figsize=(10, 6))
        plt.barh(fi_df['Feature'].head(15)[::-1], fi_df['Importance'].head(15)[::-1])
        plt.xlabel('Importance')
        plt.title('Top Features for Duplicate Detection')
        plt.tight_layout()
        results['feature_importance_chart'] = plt.gcf()
    
    # Add trade-off analysis metrics
    trade_off_data = []
    true_duplicates = len(true_duplicate_pairs)
    
    for _, row in results['thresholds_data'].iterrows():
        threshold = row['threshold']
        precision = row['precision']
        recall = row['recall']
        
        # Calculate true positives, false positives, false negatives
        tp = int(recall * true_duplicates)
        fp = int(tp * (1 - precision) / precision) if precision > 0 else 0
        fn = int(true_duplicates - tp)
        
        # Calculate cost based on business impact (example values)
        # This can be adjusted based on the specific business case
        cost_of_false_positive = 1.0  # Cost of investigating a non-duplicate
        cost_of_false_negative = 5.0  # Cost of missing a duplicate
        
        total_cost = (fp * cost_of_false_positive) + (fn * cost_of_false_negative)
        
        trade_off_data.append({
            'threshold': threshold,
            'precision': precision,
            'recall': recall,
            'true_positives': tp,
            'false_positives': fp,
            'false_negatives': fn,
            'total_cost': total_cost
        })
    
    results['trade_off_data'] = pd.DataFrame(trade_off_data)
    
    # Create ROC curve points (using AUC data if available)
    if 'auc' in results['metrics'] and results['metrics']['auc'] > 0:
        # We don't have true ROC curve data points, but we can approximate
        # from the precision-recall data we have
        
        # Create business impact chart
        plt.figure(figsize=(10, 6))
        trade_off_df = results['trade_off_data']
        plt.plot(trade_off_df['threshold'], trade_off_df['total_cost'], 'mo-')
        plt.xlabel('Threshold')
        plt.ylabel('Total Cost (Impact)')
        plt.title('Business Impact at Different Thresholds')
        plt.grid(True)
        
        # Find and mark the minimum cost threshold
        min_cost_idx = trade_off_df['total_cost'].idxmin()
        min_cost_threshold = trade_off_df.iloc[min_cost_idx]['threshold']
        min_cost = trade_off_df.iloc[min_cost_idx]['total_cost']
        
        plt.annotate(f"Optimal: {min_cost_threshold:.2f}",
                    (min_cost_threshold, min_cost),
                    textcoords="offset points",
                    xytext=(0, -20),
                    ha='center',
                    arrowprops=dict(arrowstyle="->"))
        
        plt.tight_layout()
        results['cost_impact_chart'] = plt.gcf()
    
    return results