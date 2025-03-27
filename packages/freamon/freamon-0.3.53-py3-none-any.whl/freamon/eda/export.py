"""
Export functionality for EDA reports and analysis results.

This module provides export capabilities for exploratory data analysis reports,
model performance evaluation, hyperparameter tuning results, and topic modeling
outcomes to various formats including Excel and PowerPoint.
"""
from typing import Any, Dict, List, Optional, Union, Tuple
import os
import io
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime

def export_to_excel(
    data: Dict[str, Any],
    output_path: str,
    report_type: str = "eda"
) -> None:
    """
    Export analysis results to Excel format with multiple sheets.
    
    Parameters
    ----------
    data : Dict[str, Any]
        The data to export, containing analysis results
    output_path : str
        The path where the Excel file will be saved
    report_type : str, default="eda"
        The type of report to export:
        - "eda": Exploratory data analysis
        - "model": Model performance
        - "hyperparameter": Hyperparameter tuning
        - "topic": Topic modeling
        
    Returns
    -------
    None
    """
    if not output_path.endswith('.xlsx'):
        output_path += '.xlsx'
        
    # Create Excel writer
    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        if report_type == "eda":
            _export_eda_to_excel(data, writer)
        elif report_type == "model":
            _export_model_to_excel(data, writer)
        elif report_type == "hyperparameter":
            _export_hyperparameter_to_excel(data, writer)
        elif report_type == "topic":
            _export_topic_to_excel(data, writer)
        else:
            raise ValueError(f"Unknown report type: {report_type}")
    
    print(f"Excel report saved to {output_path}")

def export_to_powerpoint(
    data: Dict[str, Any],
    output_path: str,
    report_type: str = "eda"
) -> None:
    """
    Export analysis results to PowerPoint format with multiple slides.
    
    Parameters
    ----------
    data : Dict[str, Any]
        The data to export, containing analysis results
    output_path : str
        The path where the PowerPoint file will be saved
    report_type : str, default="eda"
        The type of report to export:
        - "eda": Exploratory data analysis
        - "model": Model performance
        - "hyperparameter": Hyperparameter tuning
        - "topic": Topic modeling
    
    Returns
    -------
    None
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError(
            "python-pptx is required for PowerPoint export. "
            "Install it with: pip install python-pptx"
        )
        
    if not output_path.endswith('.pptx'):
        output_path += '.pptx'
        
    # Create presentation
    prs = Presentation()
    
    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    if report_type == "eda":
        title.text = "Exploratory Data Analysis Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        _export_eda_to_powerpoint(data, prs)
    elif report_type == "model":
        title.text = "Model Performance Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        _export_model_to_powerpoint(data, prs)
    elif report_type == "hyperparameter":
        title.text = "Hyperparameter Tuning Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        _export_hyperparameter_to_powerpoint(data, prs)
    elif report_type == "topic":
        title.text = "Topic Modeling Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        _export_topic_to_powerpoint(data, prs)
    else:
        raise ValueError(f"Unknown report type: {report_type}")
    
    # Save presentation
    prs.save(output_path)
    print(f"PowerPoint report saved to {output_path}")

def _export_eda_to_excel(data: Dict[str, Any], writer: pd.ExcelWriter) -> None:
    """Export EDA results to Excel format"""
    df = data.get('df')
    
    # Overview sheet
    overview_data = {
        'Metric': [
            'Rows', 'Columns', 'Numeric Columns', 'Categorical Columns',
            'Text Columns', 'Datetime Columns', 'Missing Values', 
            'Missing Values (%)'
        ],
        'Value': [
            len(df),
            len(df.columns),
            len(data.get('numeric_columns', [])),
            len(data.get('categorical_columns', [])),
            len(data.get('text_columns', [])),
            len(data.get('datetime_columns', [])),
            df.isna().sum().sum(),
            f"{df.isna().sum().sum() / (df.shape[0] * df.shape[1]) * 100:.2f}%"
        ]
    }
    overview_df = pd.DataFrame(overview_data)
    overview_df.to_excel(writer, sheet_name='Overview', index=False)
    
    # Sample Data sheet
    df.head(10).to_excel(writer, sheet_name='Sample Data')
    
    # Missing Values sheet
    if df.isna().sum().sum() > 0:
        missing_data = []
        for col in df.columns:
            missing_count = df[col].isna().sum()
            missing_pct = missing_count / len(df) * 100
            if missing_count > 0:
                missing_data.append({
                    'Column': col,
                    'Missing Count': missing_count,
                    'Missing %': f"{missing_pct:.2f}%"
                })
        
        if missing_data:
            missing_df = pd.DataFrame(missing_data)
            missing_df.to_excel(writer, sheet_name='Missing Values', index=False)
    
    # Numeric Columns sheet
    numeric_cols = data.get('numeric_columns', [])
    if numeric_cols:
        numeric_stats = df[numeric_cols].describe().T
        numeric_stats.to_excel(writer, sheet_name='Numeric Columns')
        
        # Correlation Matrix
        if len(numeric_cols) > 1:
            corr_matrix = df[numeric_cols].corr()
            corr_matrix.to_excel(writer, sheet_name='Correlations')
    
    # Categorical Columns sheet
    cat_cols = data.get('categorical_columns', [])
    if cat_cols:
        cat_stats = []
        for col in cat_cols:
            if col not in data.get('text_columns', []):  # Skip text columns
                value_counts = df[col].value_counts()
                unique_count = len(value_counts)
                top_value = value_counts.index[0] if not value_counts.empty else None
                top_count = value_counts.iloc[0] if not value_counts.empty else 0
                top_pct = top_count / len(df) * 100
                
                cat_stats.append({
                    'Column': col,
                    'Unique Values': unique_count,
                    'Top Value': top_value,
                    'Top Count': top_count,
                    'Top %': f"{top_pct:.2f}%"
                })
        
        if cat_stats:
            cat_df = pd.DataFrame(cat_stats)
            cat_df.to_excel(writer, sheet_name='Categorical Columns', index=False)
    
    # Text Columns sheet
    text_cols = data.get('text_columns', [])
    if text_cols:
        text_stats = []
        for col in text_cols:
            text_data = df[col].dropna().astype(str)
            char_count = text_data.str.len()
            word_count = text_data.str.split().str.len()
            
            text_stats.append({
                'Column': col,
                'Non-null Count': df[col].count(),
                'Non-null %': f"{df[col].count() / len(df) * 100:.2f}%",
                'Unique Values': df[col].nunique(),
                'Mean Char Length': f"{char_count.mean():.1f}",
                'Mean Word Count': f"{word_count.mean():.1f}",
                'Min Length': char_count.min(),
                'Max Length': char_count.max()
            })
        
        if text_stats:
            text_df = pd.DataFrame(text_stats)
            text_df.to_excel(writer, sheet_name='Text Columns', index=False)
    
    # Datetime Columns sheet
    dt_cols = data.get('datetime_columns', [])
    if dt_cols:
        dt_stats = []
        for col in dt_cols:
            try:
                dt_series = pd.to_datetime(df[col])
                dt_min = dt_series.min()
                dt_max = dt_series.max()
                dt_range = (dt_max - dt_min).days
                
                dt_stats.append({
                    'Column': col,
                    'Non-null Count': dt_series.count(),
                    'Non-null %': f"{dt_series.count() / len(df) * 100:.2f}%",
                    'Min Date': dt_min,
                    'Max Date': dt_max,
                    'Range (days)': dt_range
                })
            except:
                dt_stats.append({
                    'Column': col,
                    'Non-null Count': df[col].count(),
                    'Non-null %': f"{df[col].count() / len(df) * 100:.2f}%",
                    'Min Date': 'Error',
                    'Max Date': 'Error',
                    'Range (days)': 'Error'
                })
        
        if dt_stats:
            dt_df = pd.DataFrame(dt_stats)
            dt_df.to_excel(writer, sheet_name='Datetime Columns', index=False)
    
    # If we have multicollinearity analysis
    if 'multicollinearity' in data:
        multi_data = []
        for pair, corr in data['multicollinearity'].items():
            multi_data.append({
                'Feature 1': pair[0],
                'Feature 2': pair[1],
                'Correlation': corr
            })
        
        if multi_data:
            multi_df = pd.DataFrame(multi_data)
            multi_df.to_excel(writer, sheet_name='Multicollinearity', index=False)

def _export_eda_to_powerpoint(data: Dict[str, Any], prs) -> None:
    """Export EDA results to PowerPoint format"""
    from pptx.util import Inches, Pt
    
    df = data.get('df')
    
    # Dataset Overview slide
    overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
    overview_slide.shapes.title.text = "Dataset Overview"
    
    # Add overview table
    overview_data = [
        ('Metric', 'Value'),
        ('Rows', f"{len(df):,}"),
        ('Columns', f"{len(df.columns):,}"),
        ('Numeric Columns', f"{len(data.get('numeric_columns', [])):,}"),
        ('Categorical Columns', f"{len(data.get('categorical_columns', [])):,}"),
        ('Text Columns', f"{len(data.get('text_columns', [])):,}"),
        ('Datetime Columns', f"{len(data.get('datetime_columns', [])):,}"),
        ('Missing Values', f"{df.isna().sum().sum():,}"),
        ('Missing Values (%)', f"{df.isna().sum().sum() / (df.shape[0] * df.shape[1]) * 100:.2f}%")
    ]
    
    # Add table to slide
    top = Inches(1.5)
    left = Inches(1)
    width = Inches(8)
    height = Inches(0.8 * len(overview_data))
    
    shapes = overview_slide.shapes
    table = shapes.add_table(
        rows=len(overview_data),
        cols=2,
        left=left,
        top=top,
        width=width,
        height=height
    ).table
    
    # Set column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    
    # Add data to table
    for i, (metric, value) in enumerate(overview_data):
        table.cell(i, 0).text = str(metric)
        table.cell(i, 1).text = str(value)
    
    # Add Data Type Distribution slide if we have visualization data
    if any([
        data.get('numeric_columns', []),
        data.get('categorical_columns', []),
        data.get('text_columns', []),
        data.get('datetime_columns', [])
    ]):
        dtype_slide = prs.slides.add_slide(prs.slide_layouts[5])
        dtype_slide.shapes.title.text = "Data Type Distribution"
        
        # Create data type distribution chart
        dtype_counts = {
            'Numeric': len(data.get('numeric_columns', [])),
            'Categorical': len([c for c in data.get('categorical_columns', []) 
                               if c not in data.get('text_columns', [])]),
            'Text': len(data.get('text_columns', [])),
            'Datetime': len(data.get('datetime_columns', []))
        }
        
        # Filter out zero counts
        dtype_counts = {k: v for k, v in dtype_counts.items() if v > 0}
        
        plt.figure(figsize=(10, 6))
        bars = plt.bar(dtype_counts.keys(), dtype_counts.values(), color=sns.color_palette("viridis", len(dtype_counts)))
        plt.title('Column Data Type Distribution')
        plt.ylabel('Count')
        
        # Add count labels on top of bars
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                    f'{height:,}', ha='center')
        
        plt.tight_layout()
        
        # Save chart to temporary file
        temp_path = "temp_dtype_dist.png"
        plt.savefig(temp_path)
        plt.close()
        
        # Add chart to slide
        dtype_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
        
        # Clean up temporary file
        os.remove(temp_path)
    
    # Add Missing Values slide if there are missing values
    if df.isna().sum().sum() > 0:
        missing_slide = prs.slides.add_slide(prs.slide_layouts[5])
        missing_slide.shapes.title.text = "Missing Values Analysis"
        
        # Create missing values heatmap
        cols_with_missing = df.columns[df.isna().any()].tolist()
        
        if cols_with_missing:
            # Only proceed if there are missing values and not too many columns
            if len(cols_with_missing) <= 30:
                missing_data = df[cols_with_missing].isna()
                
                plt.figure(figsize=(10, 6))
                sns.heatmap(missing_data, cmap='viridis', cbar=False, yticklabels=False)
                plt.title('Missing Values Heatmap')
                plt.tight_layout()
                
                # Save chart to temporary file
                temp_path = "temp_missing_heatmap.png"
                plt.savefig(temp_path)
                plt.close()
                
                # Add chart to slide
                missing_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
                
                # Clean up temporary file
                os.remove(temp_path)
            
            # Add text to slide about missing values
            shapes = missing_slide.shapes
            textbox = shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
            tf = textbox.text_frame
            tf.text = f"Dataset has {df.isna().sum().sum():,} missing values ({df.isna().sum().sum() / (df.shape[0] * df.shape[1]) * 100:.2f}% of all values)"
    
    # Add Correlation slide if we have numeric correlations
    numeric_cols = data.get('numeric_columns', [])
    if len(numeric_cols) > 1:
        corr_slide = prs.slides.add_slide(prs.slide_layouts[5])
        corr_slide.shapes.title.text = "Numeric Correlations"
        
        # Create correlation heatmap
        correlation = df[numeric_cols].corr()
        mask = np.triu(np.ones_like(correlation, dtype=bool))
        
        plt.figure(figsize=(10, 8))
        sns.heatmap(correlation, mask=mask, annot=len(numeric_cols) < 10, 
                   fmt=".2f", cmap="coolwarm", vmin=-1, vmax=1, cbar=True,
                   linewidths=0.5)
        plt.title('Correlation Heatmap')
        plt.tight_layout()
        
        # Save chart to temporary file
        temp_path = "temp_correlation.png"
        plt.savefig(temp_path)
        plt.close()
        
        # Add chart to slide
        corr_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
        
        # Clean up temporary file
        os.remove(temp_path)
    
    # If we have multicollinearity analysis and enough features are highly correlated
    if 'multicollinearity' in data and len(data['multicollinearity']) > 0:
        multi_slide = prs.slides.add_slide(prs.slide_layouts[5])
        multi_slide.shapes.title.text = "Multicollinearity Analysis"
        
        # Add text about multicollinearity
        shapes = multi_slide.shapes
        textbox = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = textbox.text_frame
        tf.text = (f"Found {len(data['multicollinearity']):,} highly correlated feature pairs "
                  f"that may cause multicollinearity issues in modeling.")

def _export_model_to_excel(data: Dict[str, Any], writer: pd.ExcelWriter) -> None:
    """Export model performance results to Excel format"""
    # Overview sheet with model information
    overview_data = [
        ('Model Type', data.get('model_type', 'Unknown')),
        ('Model Library', data.get('model_library', 'Unknown')),
        ('Training Date', data.get('training_date', datetime.now().strftime('%Y-%m-%d %H:%M:%S')))
    ]
    
    # Add overall metrics
    metrics = data.get('metrics', {})
    if metrics:
        for metric_name, metric_value in metrics.items():
            if isinstance(metric_value, (int, float)):
                overview_data.append((metric_name, f"{metric_value:.4f}"))
            else:
                overview_data.append((metric_name, str(metric_value)))
    
    overview_df = pd.DataFrame(overview_data, columns=['Metric', 'Value'])
    overview_df.to_excel(writer, sheet_name='Overview', index=False)
    
    # CV Results sheet
    cv_results = data.get('cv_results', {})
    if cv_results:
        # Extract metrics (skip non-metric fields)
        cv_metrics = {}
        non_metrics = ['fold', 'train_size', 'test_size', 'train_start_date', 
                      'train_end_date', 'test_start_date', 'test_end_date',
                      'predictions', 'test_targets', 'test_dates']
        
        for metric, values in cv_results.items():
            if metric not in non_metrics and isinstance(values[0], (int, float)):
                cv_metrics[metric] = values
        
        # Create dataframe with fold-by-fold metrics
        cv_df = pd.DataFrame(cv_metrics)
        cv_df.index.name = 'Fold'
        
        # Add mean and std at the bottom
        cv_df.loc['Mean'] = cv_df.mean()
        cv_df.loc['Std'] = cv_df.std()
        
        cv_df.to_excel(writer, sheet_name='CV Results')
    
    # Feature Importance sheet
    feature_importance = data.get('feature_importance', {})
    if feature_importance:
        importance_data = []
        
        # Handle different types of feature importance
        if 'importances_mean' in feature_importance:
            # Handle scikit-learn style importance
            for feature, importance in zip(
                feature_importance.get('feature_names', []),
                feature_importance.get('importances_mean', [])
            ):
                importance_data.append({
                    'Feature': feature,
                    'Importance': importance,
                    'Std Dev': feature_importance.get('importances_std', [0] * len(feature_importance.get('importances_mean', [])))[
                        feature_importance.get('feature_names', []).index(feature)
                    ]
                })
        elif isinstance(feature_importance, dict) and len(feature_importance) > 0:
            # Handle simple dict of feature->importance mappings
            for feature, importance in feature_importance.items():
                importance_data.append({
                    'Feature': feature,
                    'Importance': importance
                })
        
        if importance_data:
            # Sort by importance descending
            importance_df = pd.DataFrame(importance_data)
            importance_df = importance_df.sort_values('Importance', ascending=False)
            importance_df.to_excel(writer, sheet_name='Feature Importance', index=False)
    
    # Prediction Results sheet (if available)
    predictions = data.get('predictions', None)
    actuals = data.get('actuals', None)
    
    if predictions is not None and actuals is not None:
        pred_df = pd.DataFrame({
            'Actual': actuals,
            'Predicted': predictions
        })
        
        # Add row index as 'Sample ID'
        pred_df.index.name = 'Sample ID'
        
        # Limit to first 1000 rows if too large
        if len(pred_df) > 1000:
            pred_df = pred_df.head(1000)
            
        pred_df.to_excel(writer, sheet_name='Predictions')
    
    # Confusion Matrix sheet (for classification)
    confusion_matrix = data.get('confusion_matrix', None)
    if confusion_matrix is not None:
        class_labels = data.get('class_labels', list(range(confusion_matrix.shape[0])))
        conf_df = pd.DataFrame(confusion_matrix, index=class_labels, columns=class_labels)
        conf_df.index.name = 'Actual'
        conf_df.columns.name = 'Predicted'
        conf_df.to_excel(writer, sheet_name='Confusion Matrix')

def _export_model_to_powerpoint(data: Dict[str, Any], prs) -> None:
    """Export model performance results to PowerPoint format"""
    from pptx.util import Inches, Pt
    
    # Model Overview slide
    overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
    overview_slide.shapes.title.text = "Model Performance Overview"
    
    # Add model info and metrics
    overview_data = [
        ('Metric', 'Value'),
        ('Model Type', data.get('model_type', 'Unknown')),
        ('Model Library', data.get('model_library', 'Unknown')),
        ('Training Date', data.get('training_date', datetime.now().strftime('%Y-%m-%d %H:%M:%S')))
    ]
    
    # Add overall metrics
    metrics = data.get('metrics', {})
    if metrics:
        for metric_name, metric_value in metrics.items():
            if isinstance(metric_value, (int, float)):
                overview_data.append((metric_name, f"{metric_value:.4f}"))
            else:
                overview_data.append((metric_name, str(metric_value)))
    
    # Add table to slide
    top = Inches(1.5)
    left = Inches(1)
    width = Inches(8)
    height = Inches(0.8 * len(overview_data))
    
    shapes = overview_slide.shapes
    table = shapes.add_table(
        rows=len(overview_data),
        cols=2,
        left=left,
        top=top,
        width=width,
        height=height
    ).table
    
    # Set column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    
    # Add data to table
    for i, (metric, value) in enumerate(overview_data):
        table.cell(i, 0).text = str(metric)
        table.cell(i, 1).text = str(value)
    
    # CV Results slide
    cv_results = data.get('cv_results', {})
    if cv_results:
        cv_slide = prs.slides.add_slide(prs.slide_layouts[5])
        cv_slide.shapes.title.text = "Cross-Validation Results"
        
        # Extract metrics (skip non-metric fields)
        metrics = {}
        non_metrics = ['fold', 'train_size', 'test_size', 'train_start_date', 
                      'train_end_date', 'test_start_date', 'test_end_date',
                      'predictions', 'test_targets', 'test_dates']
        
        for metric, values in cv_results.items():
            if metric not in non_metrics and isinstance(values[0], (int, float)):
                metrics[metric] = values
        
        # Calculate mean and std for each metric
        means = {k: np.mean(v) for k, v in metrics.items()}
        stds = {k: np.std(v) for k, v in metrics.items()}
        
        # Create bar chart with error bars
        plt.figure(figsize=(10, 6))
        x = list(means.keys())
        y = list(means.values())
        error = list(stds.values())
        
        plt.bar(x, y, yerr=error, capsize=5, color=sns.color_palette("viridis", len(x)))
        plt.title('Cross-Validation Metrics')
        plt.xticks(rotation=45, ha='right')
        plt.ylabel('Value')
        plt.tight_layout()
        
        # Save chart to temporary file
        temp_path = "temp_cv_results.png"
        plt.savefig(temp_path)
        plt.close()
        
        # Add chart to slide
        cv_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
        
        # Clean up temporary file
        os.remove(temp_path)
    
    # Feature Importance slide
    feature_importance = data.get('feature_importance', {})
    if feature_importance:
        importance_slide = prs.slides.add_slide(prs.slide_layouts[5])
        importance_slide.shapes.title.text = "Feature Importance"
        
        # Prepare feature importance data
        features = []
        importances = []
        
        # Handle different types of feature importance
        if 'importances_mean' in feature_importance:
            # Handle scikit-learn style importance
            features = feature_importance.get('feature_names', [])
            importances = feature_importance.get('importances_mean', [])
        elif isinstance(feature_importance, dict) and len(feature_importance) > 0:
            # Handle simple dict of feature->importance mappings
            for feature, importance in feature_importance.items():
                features.append(feature)
                importances.append(importance)
        
        if features and importances:
            # Sort features by importance
            sorted_idx = np.argsort(importances)
            features = [features[i] for i in sorted_idx[-20:]]  # Top 20 features
            importances = [importances[i] for i in sorted_idx[-20:]]
            
            # Create horizontal bar chart
            plt.figure(figsize=(10, 8))
            plt.barh(range(len(features)), importances, align='center')
            plt.yticks(range(len(features)), features)
            plt.xlabel('Importance')
            plt.title('Feature Importance')
            plt.tight_layout()
            
            # Save chart to temporary file
            temp_path = "temp_feature_importance.png"
            plt.savefig(temp_path)
            plt.close()
            
            # Add chart to slide
            importance_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Clean up temporary file
            os.remove(temp_path)
    
    # Add confusion matrix for classification models
    confusion_matrix = data.get('confusion_matrix', None)
    if confusion_matrix is not None:
        confusion_slide = prs.slides.add_slide(prs.slide_layouts[5])
        confusion_slide.shapes.title.text = "Confusion Matrix"
        
        class_labels = data.get('class_labels', list(range(confusion_matrix.shape[0])))
        
        plt.figure(figsize=(8, 6))
        sns.heatmap(confusion_matrix, annot=True, fmt='d', cmap='Blues',
                   xticklabels=class_labels, yticklabels=class_labels)
        plt.xlabel('Predicted')
        plt.ylabel('Actual')
        plt.title('Confusion Matrix')
        plt.tight_layout()
        
        # Save chart to temporary file
        temp_path = "temp_confusion_matrix.png"
        plt.savefig(temp_path)
        plt.close()
        
        # Add chart to slide
        confusion_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
        
        # Clean up temporary file
        os.remove(temp_path)
    
    # Add prediction scatter plot for regression models
    if 'model_type' in data and 'regression' in data['model_type'].lower():
        predictions = data.get('predictions', None)
        actuals = data.get('actuals', None)
        
        if predictions is not None and actuals is not None:
            prediction_slide = prs.slides.add_slide(prs.slide_layouts[5])
            prediction_slide.shapes.title.text = "Actual vs Predicted"
            
            plt.figure(figsize=(8, 8))
            plt.scatter(actuals, predictions, alpha=0.5)
            
            # Add perfect prediction line
            min_val = min(min(actuals), min(predictions))
            max_val = max(max(actuals), max(predictions))
            plt.plot([min_val, max_val], [min_val, max_val], 'r--')
            
            plt.xlabel('Actual')
            plt.ylabel('Predicted')
            plt.title('Actual vs Predicted Values')
            plt.axis('equal')
            plt.tight_layout()
            
            # Save chart to temporary file
            temp_path = "temp_prediction_scatter.png"
            plt.savefig(temp_path)
            plt.close()
            
            # Add chart to slide
            prediction_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Clean up temporary file
            os.remove(temp_path)

def _export_hyperparameter_to_excel(data: Dict[str, Any], writer: pd.ExcelWriter) -> None:
    """Export hyperparameter tuning results to Excel format"""
    # Overview sheet with tuning results
    overview_data = [
        ('Model Type', data.get('model_type', 'Unknown')),
        ('Optimization Algorithm', data.get('optimization_algorithm', 'Unknown')),
        ('Number of Trials', data.get('n_trials', 0)),
        ('Best Score', f"{data.get('best_score', 0):.4f}"),
        ('Optimization Metric', data.get('metric', 'Unknown')),
        ('Tuning Date', data.get('tuning_date', datetime.now().strftime('%Y-%m-%d %H:%M:%S')))
    ]
    
    overview_df = pd.DataFrame(overview_data, columns=['Metric', 'Value'])
    overview_df.to_excel(writer, sheet_name='Overview', index=False)
    
    # Best Parameters sheet
    best_params = data.get('best_params', {})
    if best_params:
        params_data = []
        for param_name, param_value in best_params.items():
            params_data.append({
                'Parameter': param_name,
                'Value': param_value
            })
        
        params_df = pd.DataFrame(params_data)
        params_df.to_excel(writer, sheet_name='Best Parameters', index=False)
    
    # Parameter Importance sheet
    param_importance = data.get('param_importance', {})
    if param_importance:
        importance_data = []
        for param_name, importance in param_importance.items():
            importance_data.append({
                'Parameter': param_name,
                'Importance': importance
            })
        
        # Sort by importance descending
        importance_df = pd.DataFrame(importance_data)
        importance_df = importance_df.sort_values('Importance', ascending=False)
        importance_df.to_excel(writer, sheet_name='Parameter Importance', index=False)
    
    # Trial History sheet
    trials = data.get('trials', [])
    if trials:
        # Format trials data
        trials_data = []
        for i, trial in enumerate(trials):
            trial_dict = {'Trial': i+1}
            
            # Add parameters
            for param_name, param_value in trial.get('params', {}).items():
                trial_dict[f"param_{param_name}"] = param_value
            
            # Add score
            trial_dict['Score'] = trial.get('value', None)
            
            trials_data.append(trial_dict)
        
        trials_df = pd.DataFrame(trials_data)
        trials_df.to_excel(writer, sheet_name='Trial History', index=False)
    
    # Parameter Ranges sheet
    param_ranges = data.get('param_ranges', {})
    if param_ranges:
        ranges_data = []
        for param_name, param_range in param_ranges.items():
            if isinstance(param_range, (list, tuple)):
                ranges_data.append({
                    'Parameter': param_name,
                    'Type': 'Categorical' if isinstance(param_range[0], str) else 'Discrete',
                    'Values': str(param_range)
                })
            elif isinstance(param_range, dict):
                if 'low' in param_range and 'high' in param_range:
                    ranges_data.append({
                        'Parameter': param_name,
                        'Type': 'Range',
                        'Min': param_range.get('low', ''),
                        'Max': param_range.get('high', ''),
                        'Log Scale': param_range.get('log', False)
                    })
                else:
                    ranges_data.append({
                        'Parameter': param_name,
                        'Type': 'Dictionary',
                        'Values': str(param_range)
                    })
        
        ranges_df = pd.DataFrame(ranges_data)
        ranges_df.to_excel(writer, sheet_name='Parameter Ranges', index=False)

def _export_hyperparameter_to_powerpoint(data: Dict[str, Any], prs) -> None:
    """Export hyperparameter tuning results to PowerPoint format"""
    from pptx.util import Inches, Pt
    
    # Tuning Overview slide
    overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
    overview_slide.shapes.title.text = "Hyperparameter Tuning Overview"
    
    # Add overview info
    overview_data = [
        ('Metric', 'Value'),
        ('Model Type', data.get('model_type', 'Unknown')),
        ('Optimization Algorithm', data.get('optimization_algorithm', 'Unknown')),
        ('Number of Trials', str(data.get('n_trials', 0))),
        ('Best Score', f"{data.get('best_score', 0):.4f}"),
        ('Optimization Metric', data.get('metric', 'Unknown')),
        ('Tuning Date', data.get('tuning_date', datetime.now().strftime('%Y-%m-%d %H:%M:%S')))
    ]
    
    # Add table to slide
    top = Inches(1.5)
    left = Inches(1)
    width = Inches(8)
    height = Inches(0.8 * len(overview_data))
    
    shapes = overview_slide.shapes
    table = shapes.add_table(
        rows=len(overview_data),
        cols=2,
        left=left,
        top=top,
        width=width,
        height=height
    ).table
    
    # Set column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    
    # Add data to table
    for i, (metric, value) in enumerate(overview_data):
        table.cell(i, 0).text = str(metric)
        table.cell(i, 1).text = str(value)
    
    # Best Parameters slide
    best_params = data.get('best_params', {})
    if best_params:
        params_slide = prs.slides.add_slide(prs.slide_layouts[5])
        params_slide.shapes.title.text = "Best Parameters"
        
        # Add parameters table
        params_data = [('Parameter', 'Value')]
        for param_name, param_value in best_params.items():
            params_data.append((param_name, str(param_value)))
        
        # Add table to slide
        top = Inches(1.5)
        left = Inches(1)
        width = Inches(8)
        height = Inches(0.8 * len(params_data))
        
        shapes = params_slide.shapes
        table = shapes.add_table(
            rows=len(params_data),
            cols=2,
            left=left,
            top=top,
            width=width,
            height=height
        ).table
        
        # Set column widths
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(4)
        
        # Add data to table
        for i, (param, value) in enumerate(params_data):
            table.cell(i, 0).text = str(param)
            table.cell(i, 1).text = str(value)
    
    # Parameter Importance slide
    param_importance = data.get('param_importance', {})
    if param_importance:
        importance_slide = prs.slides.add_slide(prs.slide_layouts[5])
        importance_slide.shapes.title.text = "Parameter Importance"
        
        # Convert to lists and sort by importance
        params = list(param_importance.keys())
        importances = list(param_importance.values())
        
        # Sort parameters by importance
        sorted_idx = np.argsort(importances)
        params = [params[i] for i in sorted_idx[-10:]]  # Top 10 parameters
        importances = [importances[i] for i in sorted_idx[-10:]]
        
        # Create horizontal bar chart
        plt.figure(figsize=(10, 6))
        plt.barh(range(len(params)), importances, align='center')
        plt.yticks(range(len(params)), params)
        plt.xlabel('Importance')
        plt.title('Parameter Importance')
        plt.tight_layout()
        
        # Save chart to temporary file
        temp_path = "temp_param_importance.png"
        plt.savefig(temp_path)
        plt.close()
        
        # Add chart to slide
        importance_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
        
        # Clean up temporary file
        os.remove(temp_path)
    
    # Trial History slide
    trials = data.get('trials', [])
    if trials:
        trial_slide = prs.slides.add_slide(prs.slide_layouts[5])
        trial_slide.shapes.title.text = "Optimization Progress"
        
        # Extract trial scores
        trial_scores = [trial.get('value', None) for trial in trials]
        
        if None not in trial_scores:
            # Create optimization progress plot
            plt.figure(figsize=(10, 6))
            plt.plot(range(1, len(trial_scores) + 1), trial_scores, 'o-')
            
            # Add best score as horizontal line
            best_score = min(trial_scores) if data.get('metric_direction', 'minimize') == 'minimize' else max(trial_scores)
            plt.axhline(y=best_score, color='r', linestyle='--', label=f'Best Score: {best_score:.4f}')
            
            plt.xlabel('Trial Number')
            plt.ylabel(data.get('metric', 'Score'))
            plt.title('Optimization Progress')
            plt.legend()
            plt.grid(True)
            plt.tight_layout()
            
            # Save chart to temporary file
            temp_path = "temp_optimization_progress.png"
            plt.savefig(temp_path)
            plt.close()
            
            # Add chart to slide
            trial_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Clean up temporary file
            os.remove(temp_path)
    
    # Parameter Interaction slide (for most important parameters)
    # This is more complex and would require additional data

def _export_topic_to_excel(data: Dict[str, Any], writer: pd.ExcelWriter) -> None:
    """Export topic modeling results to Excel format"""
    # Overview sheet with model information
    overview_data = [
        ('Topic Model Type', data.get('model_type', 'Unknown')),
        ('Number of Topics', data.get('n_topics', 0)),
        ('Vocabulary Size', data.get('vocab_size', 0)),
        ('Number of Documents', data.get('n_documents', 0)),
        ('Preprocessing', data.get('preprocessing', 'Unknown'))
    ]
    
    # Add coherence scores if available
    coherence = data.get('coherence', {})
    if coherence:
        for metric, score in coherence.items():
            overview_data.append((f"Coherence ({metric})", f"{score:.4f}"))
    
    overview_df = pd.DataFrame(overview_data, columns=['Metric', 'Value'])
    overview_df.to_excel(writer, sheet_name='Overview', index=False)
    
    # Topics sheet with term weights
    topics = data.get('topics', [])
    if topics:
        topics_data = []
        for topic_idx, terms in enumerate(topics):
            for rank, term_info in enumerate(terms[:20]):  # Top 20 terms per topic
                # Handle different term formats
                if isinstance(term_info, tuple):
                    term = term_info[0]
                    weight = term_info[1] if len(term_info) > 1 else None
                else:
                    term = term_info
                    weight = None
                
                topics_data.append({
                    'Topic': f'Topic {topic_idx + 1}',
                    'Term Rank': rank + 1,
                    'Term': term,
                    'Weight': weight
                })
        
        topics_df = pd.DataFrame(topics_data)
        topics_df.to_excel(writer, sheet_name='Topics', index=False)
    
    # Topic-Document Distribution sheet
    doc_topics = data.get('document_topics', None)
    if isinstance(doc_topics, pd.DataFrame):
        # Limit to first 1000 rows if too large
        if len(doc_topics) > 1000:
            doc_topics_out = doc_topics.head(1000)
        else:
            doc_topics_out = doc_topics.copy()
        
        doc_topics_out.to_excel(writer, sheet_name='Document Topics', index=True)
    
    # Topic Coherence sheet
    coherence_details = data.get('coherence_details', {})
    if coherence_details:
        coherence_data = []
        
        for metric, scores in coherence_details.items():
            if isinstance(scores, dict):  # Per-topic coherence
                for topic_idx, score in scores.items():
                    coherence_data.append({
                        'Metric': metric,
                        'Topic': f'Topic {topic_idx}',
                        'Score': score
                    })
            elif isinstance(scores, (int, float)):  # Overall coherence
                coherence_data.append({
                    'Metric': metric,
                    'Topic': 'Overall',
                    'Score': scores
                })
        
        if coherence_data:
            coherence_df = pd.DataFrame(coherence_data)
            coherence_df.to_excel(writer, sheet_name='Topic Coherence', index=False)
    
    # Topic-Category sheet (if categories available)
    if 'category_topics' in data:
        category_topics = data['category_topics']
        
        if isinstance(category_topics, pd.DataFrame):
            category_topics.to_excel(writer, sheet_name='Category Topics', index=True)
        elif isinstance(category_topics, dict):
            cat_topic_data = []
            
            for category, topic_dist in category_topics.items():
                if isinstance(topic_dist, dict):
                    for topic, weight in topic_dist.items():
                        cat_topic_data.append({
                            'Category': category,
                            'Topic': topic,
                            'Weight': weight
                        })
                elif isinstance(topic_dist, (list, np.ndarray)):
                    for topic_idx, weight in enumerate(topic_dist):
                        cat_topic_data.append({
                            'Category': category,
                            'Topic': f'Topic {topic_idx + 1}',
                            'Weight': weight
                        })
            
            if cat_topic_data:
                cat_topic_df = pd.DataFrame(cat_topic_data)
                cat_topic_df.to_excel(writer, sheet_name='Category Topics', index=False)
    
    # Topic Similarity sheet
    topic_similarity = data.get('topic_similarity', None)
    if topic_similarity is not None:
        if isinstance(topic_similarity, (list, np.ndarray)):
            n_topics = data.get('n_topics', len(topic_similarity))
            topic_labels = [f'Topic {i+1}' for i in range(n_topics)]
            sim_df = pd.DataFrame(topic_similarity, index=topic_labels, columns=topic_labels)
            sim_df.to_excel(writer, sheet_name='Topic Similarity')

def _export_topic_to_powerpoint(data: Dict[str, Any], prs) -> None:
    """Export topic modeling results to PowerPoint format"""
    from pptx.util import Inches, Pt
    
    # Topic Model Overview slide
    overview_slide = prs.slides.add_slide(prs.slide_layouts[1])
    overview_slide.shapes.title.text = "Topic Model Overview"
    
    # Add overview info
    overview_data = [
        ('Metric', 'Value'),
        ('Topic Model Type', data.get('model_type', 'Unknown')),
        ('Number of Topics', str(data.get('n_topics', 0))),
        ('Vocabulary Size', str(data.get('vocab_size', 0))),
        ('Number of Documents', str(data.get('n_documents', 0))),
        ('Preprocessing', data.get('preprocessing', 'Unknown'))
    ]
    
    # Add coherence scores if available
    coherence = data.get('coherence', {})
    if coherence:
        for metric, score in coherence.items():
            overview_data.append((f"Coherence ({metric})", f"{score:.4f}"))
    
    # Add table to slide
    top = Inches(1.5)
    left = Inches(1)
    width = Inches(8)
    height = Inches(0.8 * len(overview_data))
    
    shapes = overview_slide.shapes
    table = shapes.add_table(
        rows=len(overview_data),
        cols=2,
        left=left,
        top=top,
        width=width,
        height=height
    ).table
    
    # Set column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    
    # Add data to table
    for i, (metric, value) in enumerate(overview_data):
        table.cell(i, 0).text = str(metric)
        table.cell(i, 1).text = str(value)
    
    # Topic Word Charts slide for each topic
    topics = data.get('topics', [])
    if topics:
        # Limit to top N topics to avoid too many slides
        max_topic_slides = min(len(topics), 5)
        
        for topic_idx in range(max_topic_slides):
            terms = topics[topic_idx]
            
            # Create slide for this topic
            topic_slide = prs.slides.add_slide(prs.slide_layouts[5])
            topic_slide.shapes.title.text = f"Topic {topic_idx + 1}: Top Terms"
            
            # Extract term data
            if isinstance(terms[0], tuple):
                term_words = [term[0] for term in terms[:15]]
                term_weights = [term[1] for term in terms[:15]]
            else:
                term_words = terms[:15]
                term_weights = [1.0] * len(term_words)  # Default weights
            
            # Create horizontal bar chart for term weights
            plt.figure(figsize=(10, 6))
            y_pos = range(len(term_words))
            plt.barh(y_pos, term_weights, align='center')
            plt.yticks(y_pos, term_words)
            plt.xlabel('Weight')
            plt.title(f'Top Terms for Topic {topic_idx + 1}')
            plt.tight_layout()
            
            # Save chart to temporary file
            temp_path = f"temp_topic_{topic_idx}.png"
            plt.savefig(temp_path)
            plt.close()
            
            # Add chart to slide
            topic_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Clean up temporary file
            os.remove(temp_path)
    
    # Topic Distribution slide
    if 'document_topics' in data:
        doc_topics = data['document_topics']
        
        if isinstance(doc_topics, pd.DataFrame):
            dist_slide = prs.slides.add_slide(prs.slide_layouts[5])
            dist_slide.shapes.title.text = "Topic Distribution"
            
            # Create topic distribution bar chart
            mean_topic_dist = doc_topics.mean()
            
            plt.figure(figsize=(10, 6))
            plt.bar(range(len(mean_topic_dist)), mean_topic_dist, color=sns.color_palette("viridis", len(mean_topic_dist)))
            plt.xticks(range(len(mean_topic_dist)), mean_topic_dist.index)
            plt.xlabel('Topic')
            plt.ylabel('Average Weight')
            plt.title('Average Topic Distribution Across Documents')
            plt.tight_layout()
            
            # Save chart to temporary file
            temp_path = "temp_topic_dist.png"
            plt.savefig(temp_path)
            plt.close()
            
            # Add chart to slide
            dist_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Clean up temporary file
            os.remove(temp_path)
    
    # Category-Topic slide (if categories available)
    if 'category_topics' in data:
        category_topics = data['category_topics']
        
        if isinstance(category_topics, pd.DataFrame):
            cat_slide = prs.slides.add_slide(prs.slide_layouts[5])
            cat_slide.shapes.title.text = "Topic Distribution by Category"
            
            # Create heatmap of category-topic relationships
            plt.figure(figsize=(12, 8))
            sns.heatmap(category_topics, cmap='viridis', annot=False, cbar=True)
            plt.xlabel('Topics')
            plt.ylabel('Categories')
            plt.title('Topic Distribution by Category')
            plt.tight_layout()
            
            # Save chart to temporary file
            temp_path = "temp_category_topic.png"
            plt.savefig(temp_path)
            plt.close()
            
            # Add chart to slide
            cat_slide.shapes.add_picture(temp_path, Inches(0.5), Inches(1.5), width=Inches(9))
            
            # Clean up temporary file
            os.remove(temp_path)
    
    # Topic Similarity slide if available
    topic_similarity = data.get('topic_similarity', None)
    if topic_similarity is not None:
        sim_slide = prs.slides.add_slide(prs.slide_layouts[5])
        sim_slide.shapes.title.text = "Topic Similarity Matrix"
        
        # Create topic similarity heatmap
        plt.figure(figsize=(10, 8))
        sns.heatmap(topic_similarity, cmap='coolwarm', annot=True, fmt=".2f",
                    xticklabels=[f'T{i+1}' for i in range(len(topic_similarity))],
                    yticklabels=[f'T{i+1}' for i in range(len(topic_similarity))])
        plt.title('Topic Similarity Matrix')
        plt.tight_layout()
        
        # Save chart to temporary file
        temp_path = "temp_topic_similarity.png"
        plt.savefig(temp_path)
        plt.close()
        
        # Add chart to slide
        sim_slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), width=Inches(8))
        
        # Clean up temporary file
        os.remove(temp_path)